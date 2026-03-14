"""
generate_slides.py — Build a PowerPoint deck with YouTube analysis charts and insights

No Google OAuth, no credentials.json, no API calls.
Generates a local .pptx file using python-pptx.

Usage:
    python tools/generate_slides.py \
        --analysis .tmp/analysis_results.json \
        --charts-dir .tmp/charts/ \
        --output .tmp/youtube_trends_report.pptx \
        --title "AI & Automation YouTube Trends"
"""

import argparse
import json
import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Design constants
DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0xE9, 0x45, 0x60)
BLUE = RGBColor(0x0F, 0x34, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF6, 0xFA)
TEXT_COLOR = RGBColor(0x33, 0x33, 0x33)
SUBTLE = RGBColor(0x88, 0x88, 0x88)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    """Set slide background to a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, text, left, top, width, height,
                 font_size=14, color=TEXT_COLOR, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a styled text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_multiline_text(slide, lines, left, top, width, height,
                       font_size=14, color=TEXT_COLOR, line_spacing=1.5):
    """Add multiple lines of text to a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(font_size * 0.5)
    return tf


def add_metric_box(slide, value, label, left, top, width, height,
                   value_color=DARK, label_color=SUBTLE):
    """Add a metric display (big number + small label)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = str(value)
    p.font.size = Pt(36)
    p.font.color.rgb = value_color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(11)
    p2.font.color.rgb = label_color
    p2.alignment = PP_ALIGN.CENTER


def build_title_slide(prs, analysis):
    """Slide 1: Title slide with dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide, DARK)

    summary = analysis.get("summary", {})
    date_str = summary.get("fetch_date", datetime.now().strftime("%Y-%m-%d"))

    add_text_box(slide, "AI & Automation\nYouTube Trends",
                 Inches(0.8), Inches(1.5), Inches(11), Inches(2),
                 font_size=44, color=WHITE, bold=True)

    add_text_box(slide, f"Analysis Report — {date_str}",
                 Inches(0.8), Inches(3.8), Inches(11), Inches(0.5),
                 font_size=18, color=ACCENT)

    total_v = summary.get("total_videos", 0)
    total_c = summary.get("total_channels", 0)
    days = summary.get("days_back", 30)
    add_text_box(slide, f"{total_v} videos  ·  {total_c} channels  ·  Last {days} days",
                 Inches(0.8), Inches(4.5), Inches(11), Inches(0.4),
                 font_size=13, color=LIGHT_GRAY)


def build_metrics_slide(prs, analysis):
    """Slide 2: Key metrics dashboard."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    summary = analysis.get("summary", {})

    add_text_box(slide, "Key Metrics",
                 Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 font_size=32, color=DARK, bold=True)

    # Metric boxes in a row
    metrics = [
        (f"{summary.get('total_videos', 0):,}", "Videos Analyzed"),
        (f"{summary.get('total_channels', 0):,}", "Channels"),
        (f"{summary.get('avg_views', 0):,}", "Avg Views"),
        (f"{summary.get('avg_engagement_rate', 0)}%", "Avg Engagement"),
    ]

    box_width = Inches(2.8)
    start_x = Inches(0.5)
    for i, (value, label) in enumerate(metrics):
        left = start_x + (box_width + Inches(0.2)) * i
        add_metric_box(slide, value, label, left, Inches(1.8), box_width, Inches(1.5),
                       value_color=ACCENT if i == 3 else DARK)

    # Additional stats below
    lines = [
        f"Median Views:  {summary.get('median_views', 0):,}",
        f"Top Video Views:  {summary.get('max_views', 0):,}",
        f"Median Engagement:  {summary.get('median_engagement_rate', 0)}%",
        f"Keywords Tracked:  {', '.join(summary.get('keywords', [])[:4])}",
    ]
    add_multiline_text(slide, lines, Inches(0.8), Inches(4.0), Inches(11), Inches(3),
                       font_size=14, color=TEXT_COLOR)


def build_chart_slide(prs, title, chart_path, subtitle=None):
    """Generic slide with a title and a chart image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, title,
                 Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
                 font_size=28, color=DARK, bold=True)

    if chart_path and Path(chart_path).exists():
        slide.shapes.add_picture(
            str(chart_path),
            Inches(0.5), Inches(1.0),
            Inches(12.3), Inches(5.8),
        )
    else:
        add_text_box(slide, "Chart not available",
                     Inches(3), Inches(3), Inches(7), Inches(1),
                     font_size=18, color=SUBTLE, alignment=PP_ALIGN.CENTER)

    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.8), Inches(6.9), Inches(11), Inches(0.4),
                     font_size=11, color=BLUE)


def build_rising_channels_slide(prs, rising_channels):
    """Slide: Rising channels list."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, "Rising Channels",
                 Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 font_size=32, color=DARK, bold=True)

    lines = []
    for i, ch in enumerate(rising_channels[:7], 1):
        subs = ch.get("subscriber_count", 0)
        subs_str = f"{subs:,}" if subs < 1_000_000 else f"{subs / 1_000_000:.1f}M"
        vpv = ch.get("views_per_video", 0)
        lines.append(f"{i}.  {ch.get('title', 'Unknown')}")
        lines.append(f"     {subs_str} subscribers  ·  {vpv:,} views/video")
        lines.append("")

    if not lines:
        lines = ["No channel data available."]

    add_multiline_text(slide, lines, Inches(0.8), Inches(1.3), Inches(11), Inches(5.5),
                       font_size=15, color=TEXT_COLOR)


def build_insights_slide(prs, insights):
    """Slide: Key insights with dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)

    add_text_box(slide, "Key Insights",
                 Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 font_size=32, color=ACCENT, bold=True)

    bullet_lines = [f"→  {insight}" for insight in insights[:7]]
    if not bullet_lines:
        bullet_lines = ["→  No insights generated yet."]

    add_multiline_text(slide, bullet_lines, Inches(0.8), Inches(1.5), Inches(11), Inches(5),
                       font_size=16, color=WHITE, line_spacing=1.8)


def main():
    parser = argparse.ArgumentParser(description="Generate PowerPoint YouTube analysis deck")
    parser.add_argument("--analysis", required=True, help="Path to analysis results JSON")
    parser.add_argument("--charts-dir", default=".tmp/charts/", help="Directory containing chart PNGs")
    parser.add_argument("--output", default=".tmp/youtube_trends_report.pptx", help="Output .pptx path")
    parser.add_argument("--title", default="AI & Automation YouTube Trends", help="Presentation title")
    args = parser.parse_args()

    analysis_path = Path(args.analysis)
    if not analysis_path.exists():
        print(f"ERROR: Analysis file not found: {args.analysis}", file=sys.stderr)
        sys.exit(1)

    analysis = json.loads(analysis_path.read_text(encoding="utf-8"))
    charts_dir = Path(args.charts_dir)

    print("Generating PowerPoint deck...")

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Build slides
    build_title_slide(prs, analysis)
    print("  [1/10] Title slide")

    build_metrics_slide(prs, analysis)
    print("  [2/10] Key metrics")

    # Chart slides
    chart_slides = [
        ("Top Videos by Views", "top_videos.png", None),
        ("Trending Keywords & Topics", "trending_keywords.png", None),
        ("Engagement Analysis", "engagement_dist.png", None),
        ("Views vs Engagement Rate", "views_engagement.png", None),
        ("Upload Patterns", "upload_heatmap.png", None),
        ("Performance by Video Duration", "duration_perf.png", None),
    ]

    # Add subtitle for upload patterns
    peaks = analysis.get("upload_patterns", {}).get("peak_windows", [])
    if peaks:
        chart_slides[4] = (
            "Upload Patterns",
            "upload_heatmap.png",
            f"Peak upload time: {peaks[0]['day']}s at {peaks[0]['hour']}:00 UTC"
        )

    for i, (title, filename, subtitle) in enumerate(chart_slides, 3):
        chart_path = charts_dir / filename
        build_chart_slide(prs, title, str(chart_path), subtitle)
        status = "OK" if chart_path.exists() else "no chart"
        print(f"  [{i}/10] {title} ({status})")

    build_rising_channels_slide(prs, analysis.get("rising_channels", []))
    print("  [9/10] Rising channels")

    build_insights_slide(prs, analysis.get("insights", []))
    print("  [10/10] Key insights")

    # Save
    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    print(f"\nDone! Saved to: {output_path}")
    print(json.dumps({
        "status": "success",
        "output": str(output_path),
        "slides_count": len(prs.slides),
    }, indent=2))


if __name__ == "__main__":
    main()
